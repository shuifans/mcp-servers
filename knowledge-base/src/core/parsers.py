from dataclasses import dataclass, field
from pathlib import Path


SUPPORTED = {".pdf", ".docx", ".pptx", ".xlsx", ".md", ".txt", ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".go", ".rs", ".json", ".yaml", ".yml"}


@dataclass
class ParsedPart:
    text: str
    location: dict = field(default_factory=dict)


def parse_file(path: Path) -> list[ParsedPart]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        import fitz
        with fitz.open(path) as doc:
            return [ParsedPart(page.get_text(), {"page": i + 1}) for i, page in enumerate(doc) if page.get_text().strip()]
    if suffix == ".docx":
        from docx import Document
        doc = Document(path)
        parts, section, buf = [], None, []
        for p in doc.paragraphs:
            text = p.text.strip()
            if not text:
                continue
            if p.style and p.style.name.startswith("Heading"):
                if buf:
                    parts.append(ParsedPart("\n".join(buf), {"section": section}))
                section, buf = text, []
            else:
                buf.append(text)
        if buf:
            parts.append(ParsedPart("\n".join(buf), {"section": section}))
        return parts
    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        return [ParsedPart("\n".join(s.text for shape in slide.shapes if hasattr(shape, "text_frame") for s in shape.text_frame.paragraphs if s.text.strip()), {"slide": i + 1}) for i, slide in enumerate(prs.slides)]
    if suffix == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            rows = ["\t".join("" if c is None else str(c) for c in row) for row in ws.iter_rows(values_only=True)]
            text = "\n".join(r for r in rows if r.strip())
            if text:
                parts.append(ParsedPart(text, {"sheet": ws.title, "cell_range": ws.calculate_dimension()}))
        return parts
    return [ParsedPart(path.read_text(encoding="utf-8", errors="ignore"), {})]


def chunk_parts(parts: list[ParsedPart], size: int = 1000, overlap: int = 120) -> list[ParsedPart]:
    chunks = []
    for part in parts:
        text = " ".join(part.text.split())
        start = 0
        while start < len(text):
            end = min(start + size, len(text))
            chunks.append(ParsedPart(text[start:end], part.location))
            if end == len(text):
                break
            start = end - overlap
    return chunks

